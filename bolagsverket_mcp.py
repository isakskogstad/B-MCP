#!/usr/bin/env python3
"""
Bolagsverket MCP Server v3.0
============================
Hämtar och analyserar företagsdata från Bolagsverkets API "Värdefulla datamängder".

MCP Best Practices Implementation:
- Tools: Granulära verktyg för specifika operationer
- Resources: Passiv data via URI-schema
- Prompts: Fördefinierade arbetsflöden
- Logging: stderr (aldrig stdout för STDIO-transport)
- Strukturerad felhantering
"""

import json
import uuid
import httpx
import zipfile
import sys
import os
import re
import csv
import logging
from datetime import datetime, timedelta
from typing import Optional, Dict, Any, List, Tuple
from enum import Enum
from io import BytesIO, StringIO
from dataclasses import dataclass, asdict, field
from pathlib import Path

from bs4 import BeautifulSoup
from mcp.server.fastmcp import FastMCP
from pydantic import BaseModel, Field, ConfigDict

# =============================================================================
# FÖRBÄTTRING #1: Logging till stderr (ALDRIG stdout för STDIO-transport)
# =============================================================================
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    stream=sys.stderr  # KRITISKT: Måste vara stderr för MCP STDIO
)
logger = logging.getLogger("bolagsverket_mcp")

# Valfria imports för export
try:
    import openpyxl
    from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    logger.warning("openpyxl ej installerat - Excel-export inaktiverad")

try:
    from weasyprint import HTML, CSS
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("weasyprint ej installerat - PDF-export inaktiverad")

try:
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx ej installerat - Word-export inaktiverad")

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx ej installerat - PowerPoint-export inaktiverad")


# =============================================================================
# Konfiguration (läses från miljövariabler)
# =============================================================================

def _get_required_env(name: str) -> str:
    """Hämta obligatorisk miljövariabel."""
    value = os.environ.get(name)
    if not value:
        raise ValueError(
            f"Miljövariabel {name} saknas. "
            f"Sätt den via: export {name}='ditt-värde' "
            f"eller skapa en .env-fil."
        )
    return value


# OAuth2 credentials - MÅSTE sättas via miljövariabler
# Skapa konto på https://portal.api.bolagsverket.se för att få credentials
CLIENT_ID = os.environ.get("BOLAGSVERKET_CLIENT_ID", "")
CLIENT_SECRET = os.environ.get("BOLAGSVERKET_CLIENT_SECRET", "")

# API endpoints (publika, behöver ej skyddas)
TOKEN_URL = "https://portal.api.bolagsverket.se/oauth2/token"
BASE_URL = "https://gw.api.bolagsverket.se/vardefulla-datamangder/v1"
SCOPE = "vardefulla-datamangder:read vardefulla-datamangder:ping"

OUTPUT_DIR = Path.home() / "Downloads" / "bolagsverket"


def _validate_credentials():
    """Validera att credentials är konfigurerade."""
    if not CLIENT_ID or not CLIENT_SECRET:
        logger.error("BOLAGSVERKET_CLIENT_ID och BOLAGSVERKET_CLIENT_SECRET måste sättas!")
        raise ValueError(
            "API-credentials saknas. Sätt miljövariablerna:\n"
            "  export BOLAGSVERKET_CLIENT_ID='din-client-id'\n"
            "  export BOLAGSVERKET_CLIENT_SECRET='din-client-secret'\n\n"
            "Skapa credentials på: https://portal.api.bolagsverket.se"
        )


# =============================================================================
# FÖRBÄTTRING #10: Strukturerade felkoder
# =============================================================================

class ErrorCode(str, Enum):
    COMPANY_NOT_FOUND = "COMPANY_NOT_FOUND"
    ANNUAL_REPORT_NOT_FOUND = "ANNUAL_REPORT_NOT_FOUND"
    API_ERROR = "API_ERROR"
    AUTH_ERROR = "AUTH_ERROR"
    PARSE_ERROR = "PARSE_ERROR"
    INVALID_INPUT = "INVALID_INPUT"
    EXPORT_ERROR = "EXPORT_ERROR"


@dataclass
class MCPError:
    """Strukturerat fel enligt MCP best practices."""
    code: ErrorCode
    message: str
    details: Dict[str, Any] = field(default_factory=dict)
    
    def to_response(self) -> str:
        """Formatera som MCP-kompatibelt felmeddelande."""
        return json.dumps({
            "isError": True,
            "errorCode": self.code.value,
            "message": self.message,
            "details": self.details
        }, ensure_ascii=False, indent=2)


def handle_error(code: ErrorCode, message: str, **details) -> str:
    """Skapa strukturerat felmeddelande."""
    error = MCPError(code=code, message=message, details=details)
    logger.error(f"{code.value}: {message} - {details}")
    return error.to_response()


# =============================================================================
# Dataklasser
# =============================================================================

@dataclass
class Person:
    fornamn: str
    efternamn: str
    roll: str
    datum: Optional[str] = None
    
    @property
    def fullnamn(self) -> str:
        return f"{self.fornamn} {self.efternamn}".strip()


@dataclass
class Nyckeltal:
    nettoomsattning: Optional[int] = None
    resultat_efter_finansiella: Optional[int] = None
    arets_resultat: Optional[int] = None
    eget_kapital: Optional[int] = None
    balansomslutning: Optional[int] = None
    soliditet: Optional[float] = None
    antal_anstallda: Optional[int] = None
    vinstmarginal: Optional[float] = None
    roe: Optional[float] = None
    
    def berakna_nyckeltal(self):
        if self.nettoomsattning and self.arets_resultat:
            self.vinstmarginal = round((self.arets_resultat / self.nettoomsattning) * 100, 2)
        if self.eget_kapital and self.arets_resultat and self.eget_kapital > 0:
            self.roe = round((self.arets_resultat / self.eget_kapital) * 100, 2)


@dataclass
class Arsredovisning:
    org_nummer: str
    foretag_namn: str
    rakenskapsar_start: str
    rakenskapsar_slut: str
    nyckeltal: Nyckeltal
    personer: List[Person]
    balansrakning: Dict[str, Any]
    resultatrakning: Dict[str, Any]
    noter: Dict[str, str]
    metadata: Dict[str, str]


@dataclass 
class CompanyInfo:
    """Grundläggande företagsinformation."""
    org_nummer: str
    namn: str
    organisationsform: str
    juridisk_form: Optional[str]
    registreringsdatum: str
    status: str
    avregistreringsdatum: Optional[str]
    adress: Dict[str, str]
    verksamhet: Optional[str]
    sni_koder: List[Dict[str, str]]
    sate: Optional[str]


# =============================================================================
# Token-hantering
# =============================================================================

class TokenManager:
    def __init__(self):
        self.access_token: Optional[str] = None
        self.token_expiry: Optional[datetime] = None

    def get_token(self, force_refresh: bool = False) -> str:
        # Validera credentials innan API-anrop
        _validate_credentials()

        if not force_refresh and self.access_token and self.token_expiry:
            if datetime.now() < self.token_expiry:
                return self.access_token

        logger.info("Hämtar ny OAuth2-token...")
        
        with httpx.Client(timeout=30.0) as client:
            response = client.post(
                TOKEN_URL,
                headers={"Content-Type": "application/x-www-form-urlencoded"},
                data={
                    "grant_type": "client_credentials",
                    "client_id": CLIENT_ID,
                    "client_secret": CLIENT_SECRET,
                    "scope": SCOPE
                }
            )
        
        if response.status_code != 200:
            logger.error(f"Token-fel: {response.status_code}")
            raise Exception(f"Token-fel: {response.status_code} - {response.text}")
        
        data = response.json()
        self.access_token = data["access_token"]
        expires_in = data.get("expires_in", 3600)
        self.token_expiry = datetime.now() + timedelta(seconds=expires_in - 60)
        
        logger.info("Token hämtad, giltig i %d sekunder", expires_in)
        return self.access_token


token_manager = TokenManager()


# =============================================================================
# FÖRBÄTTRING #7: ServerCapabilities
# =============================================================================

mcp = FastMCP("bolagsverket")


# =============================================================================
# API-hjälpfunktioner
# =============================================================================

def clean_org_nummer(org_nummer: str) -> str:
    return org_nummer.replace("-", "").replace(" ", "")


def format_org_nummer(org_nummer: str) -> str:
    clean = clean_org_nummer(org_nummer)
    if len(clean) == 10:
        return f"{clean[:6]}-{clean[6:]}"
    return clean


def validate_org_nummer(org_nummer: str) -> Tuple[bool, str]:
    """Validera organisationsnummer."""
    clean = clean_org_nummer(org_nummer)
    if not clean.isdigit():
        return False, "Organisationsnummer får endast innehålla siffror"
    if len(clean) not in (10, 12):
        return False, "Organisationsnummer måste vara 10 eller 12 siffror"
    return True, clean


def make_api_request(method: str, endpoint: str, json_body: Optional[Dict] = None) -> Dict[str, Any]:
    token = token_manager.get_token()
    
    headers = {
        "Authorization": f"Bearer {token}",
        "X-Request-Id": str(uuid.uuid4())
    }
    
    if json_body:
        headers["Content-Type"] = "application/json"
    
    url = f"{BASE_URL}{endpoint}"
    logger.debug(f"API-anrop: {method} {endpoint}")
    
    with httpx.Client(timeout=30.0) as client:
        if method == "GET":
            response = client.get(url, headers=headers)
        elif method == "POST":
            response = client.post(url, headers=headers, json=json_body)
        else:
            raise ValueError(f"Okänd HTTP-metod: {method}")
    
    if response.status_code != 200:
        logger.error(f"API-fel: {response.status_code} - {response.text[:200]}")
        try:
            error_data = response.json()
            raise Exception(error_data.get("detail", f"HTTP {response.status_code}"))
        except json.JSONDecodeError:
            raise Exception(f"HTTP {response.status_code}: {response.text[:200]}")
    
    return response.json()


def download_document_bytes(dokument_id: str) -> bytes:
    token = token_manager.get_token()
    
    headers = {
        "Authorization": f"Bearer {token}",
        "Accept": "application/zip",
        "X-Request-Id": str(uuid.uuid4())
    }
    
    url = f"{BASE_URL}/dokument/{dokument_id}"
    logger.info(f"Laddar ner dokument: {dokument_id}")
    
    with httpx.Client(timeout=60.0) as client:
        response = client.get(url, headers=headers)
    
    if response.status_code != 200:
        raise Exception(f"HTTP {response.status_code}: Kunde inte ladda ner dokument")
    
    return response.content


def ensure_output_dir() -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return OUTPUT_DIR


def fetch_company_info(org_nummer: str) -> CompanyInfo:
    """Hämta och strukturera företagsinformation."""
    clean_nr = clean_org_nummer(org_nummer)
    data = make_api_request("POST", "/organisationer", {"identitetsbeteckning": clean_nr})
    
    orgs = data.get("organisationer", [])
    if not orgs:
        raise Exception(f"Företaget {org_nummer} hittades inte")
    
    org = orgs[0]
    
    # Extrahera data
    namn = org.get("organisationsnamn", {}).get("organisationsnamnLista", [{}])[0].get("namn", "Okänt")
    
    avreg = org.get("avregistreradOrganisation", {})
    status = "Avregistrerad" if avreg and avreg.get("avregistreringsdatum") else "Aktiv"
    
    adress_data = org.get("postadressOrganisation", {}).get("postadress", {})
    adress = {
        "utdelningsadress": adress_data.get("utdelningsadress", ""),
        "postnummer": adress_data.get("postnummer", ""),
        "postort": adress_data.get("postort", "")
    }
    
    sni = org.get("naringsgrenOrganisation", {}).get("sni", [])
    sni_koder = [{"kod": s.get("kod", ""), "klartext": s.get("klartext", "")} for s in sni if s.get("kod")]
    
    return CompanyInfo(
        org_nummer=format_org_nummer(clean_nr),
        namn=namn,
        organisationsform=org.get("organisationsform", {}).get("klartext", "-"),
        juridisk_form=org.get("juridiskForm", {}).get("klartext"),
        registreringsdatum=org.get("organisationsdatum", {}).get("registreringsdatum", "-"),
        status=status,
        avregistreringsdatum=avreg.get("avregistreringsdatum") if avreg else None,
        adress=adress,
        verksamhet=org.get("verksamhetsbeskrivning", {}).get("beskrivning"),
        sni_koder=sni_koder,
        sate=org.get("sate", {}).get("lan")
    )


# =============================================================================
# iXBRL Parser
# =============================================================================

class IXBRLParser:
    """Parser för iXBRL (Inline XBRL) årsredovisningar."""
    
    def __init__(self, xhtml_content: str):
        self.soup = BeautifulSoup(xhtml_content, 'lxml')
        self._cache = {}
    
    def _get_value(self, name_pattern: str, context: str = None, numeric: bool = True) -> Optional[Any]:
        tag_type = 'ix:nonfraction' if numeric else 'ix:nonnumeric'
        
        def name_match(x):
            return x and name_pattern.lower() in x.lower()
        
        attrs = {'name': name_match}
        if context:
            attrs['contextref'] = context
        
        tag = self.soup.find(tag_type, attrs)
        if tag:
            value = tag.text.strip()
            if numeric:
                value = value.replace(' ', '').replace(',', '.').replace('−', '-')
                try:
                    scale = int(tag.get('scale', '0'))
                    return int(float(value) * (10 ** scale))
                except ValueError:
                    return None
            return value
        return None
    
    def get_metadata(self) -> Dict[str, str]:
        return {
            'foretag_namn': self._get_value('ForetagetsNamn', numeric=False) or '',
            'org_nummer': self._get_value('Organisationsnummer', numeric=False) or '',
            'rakenskapsar_start': self._get_value('RakenskapsarForstaDag', numeric=False) or '',
            'rakenskapsar_slut': self._get_value('RakenskapsarSistaDag', numeric=False) or '',
            'undertecknat_datum': self._get_value('UndertecknandeDatum', numeric=False) or '',
            'sate': self._get_value('ForetagetsSate', numeric=False) or '',
        }
    
    def get_nyckeltal(self, period: str = 'period0') -> Nyckeltal:
        balans = period.replace('period', 'balans')
        
        nyckeltal = Nyckeltal(
            nettoomsattning=self._get_value('Nettoomsattning', period),
            resultat_efter_finansiella=self._get_value('ResultatEfterFinansiellaPoster', period),
            arets_resultat=self._get_value('AretsResultat', period),
            eget_kapital=self._get_value('EgetKapital', balans),
            balansomslutning=self._get_value('Tillgangar', balans) or self._get_value('SummaEgetKapitalSkulder', balans),
            soliditet=self._get_value('Soliditet', balans),
            antal_anstallda=self._get_value('MedelantalAnstallda', period),
        )
        
        nyckeltal.berakna_nyckeltal()
        return nyckeltal
    
    def get_flerarsoversikt(self) -> Dict[str, Nyckeltal]:
        oversikt = {}
        for i in range(4):
            period = f'period{i}'
            nyckeltal = self.get_nyckeltal(period)
            if nyckeltal.nettoomsattning is not None:
                oversikt[period] = nyckeltal
        return oversikt
    
    def get_personer(self) -> List[Person]:
        personer = []
        seen = set()
        
        patterns = [
            ('UnderskriftFaststallelseintygForetradareTilltalsnamn', 
             'UnderskriftFaststallelseintygForetradareEfternamn',
             'UnderskriftFaststallelseintygForetradareForetradarroll'),
            ('UnderskriftHandlingTilltalsnamn', 'UnderskriftHandlingEfternamn', None),
            ('UnderskriftRevisionsberattelseRevisorTilltalsnamn', 
             'UnderskriftRevisionsberattelseRevisorEfternamn',
             'UnderskriftRevisionsberattelseRevisorTitel'),
        ]
        
        for fornamn_pat, efternamn_pat, roll_pat in patterns:
            for tag in self.soup.find_all('ix:nonnumeric', {'name': lambda x: x and fornamn_pat in x}):
                fornamn = tag.text.strip()
                
                tuple_ref = tag.get('tupleref')
                efternamn = ''
                roll = ''
                
                if tuple_ref:
                    efternamn_tag = self.soup.find('ix:nonnumeric', {
                        'name': lambda x: x and efternamn_pat in x,
                        'tupleref': tuple_ref
                    })
                    if efternamn_tag:
                        efternamn = efternamn_tag.text.strip()
                    
                    if roll_pat:
                        roll_tag = self.soup.find('ix:nonnumeric', {
                            'name': lambda x: x and roll_pat in x,
                            'tupleref': tuple_ref
                        })
                        if roll_tag:
                            roll = roll_tag.text.strip()
                
                if not roll:
                    if 'Revisor' in fornamn_pat:
                        roll = 'Revisor'
                    elif 'Foretradar' in fornamn_pat:
                        roll = 'Företrädare'
                    else:
                        roll = 'Styrelseledamot'
                
                key = (fornamn, efternamn, roll)
                if key not in seen and fornamn:
                    seen.add(key)
                    personer.append(Person(fornamn=fornamn, efternamn=efternamn, roll=roll))
        
        return personer
    
    def get_balansrakning(self, period: str = 'balans0') -> Dict[str, Any]:
        return {
            'tillgangar': {
                'immateriella': self._get_value('ImmateriellAnlaggningstillgangar', period),
                'materiella': self._get_value('MateriellaAnlaggningstillgangar', period),
                'finansiella': self._get_value('FinansiellaAnlaggningstillgangar', period),
                'varulager': self._get_value('VarulagerMm', period),
                'kundfordringar': self._get_value('Kundfordringar', period),
                'kassa_bank': self._get_value('KassaBank', period),
                'summa_omsattning': self._get_value('Omsattningstillgangar', period),
                'summa_tillgangar': self._get_value('Tillgangar', period),
            },
            'eget_kapital_skulder': {
                'aktiekapital': self._get_value('Aktiekapital', period),
                'balanserat_resultat': self._get_value('BalanseratResultat', period),
                'arets_resultat': self._get_value('AretsResultatEgetKapital', period),
                'summa_eget_kapital': self._get_value('EgetKapital', period),
                'langfristiga_skulder': self._get_value('LangfristigaSkulder', period),
                'kortfristiga_skulder': self._get_value('KortfristigaSkulder', period),
                'leverantorsskulder': self._get_value('Leverantorsskulder', period),
                'summa_skulder': self._get_value('Skulder', period),
            }
        }
    
    def get_resultatrakning(self, period: str = 'period0') -> Dict[str, Any]:
        return {
            'nettoomsattning': self._get_value('Nettoomsattning', period),
            'ovriga_rorelseinktakter': self._get_value('OvrigaRorelseintakter', period),
            'summa_intakter': self._get_value('RorelseintakterLagerforandringarMm', period),
            'varor_handelsvaror': self._get_value('HandelsvarorKostnader', period),
            'ovriga_externa_kostnader': self._get_value('OvrigaExternaKostnader', period),
            'personalkostnader': self._get_value('Personalkostnader', period),
            'avskrivningar': self._get_value('AvskrivningarNedskrivningarMateriellaImmateriellaAnlaggningstillgangar', period),
            'rorelseresultat': self._get_value('Rorelseresultat', period),
            'finansiella_intakter': self._get_value('FinansiellaIntakter', period),
            'finansiella_kostnader': self._get_value('FinansiellaKostnader', period),
            'resultat_efter_finansiella': self._get_value('ResultatEfterFinansiellaPoster', period),
            'skatt': self._get_value('SkattAretsResultat', period),
            'arets_resultat': self._get_value('AretsResultat', period),
        }
    
    def parse_full(self) -> Arsredovisning:
        metadata = self.get_metadata()
        
        return Arsredovisning(
            org_nummer=metadata['org_nummer'],
            foretag_namn=metadata['foretag_namn'],
            rakenskapsar_start=metadata['rakenskapsar_start'],
            rakenskapsar_slut=metadata['rakenskapsar_slut'],
            nyckeltal=self.get_nyckeltal(),
            personer=self.get_personer(),
            balansrakning=self.get_balansrakning(),
            resultatrakning=self.get_resultatrakning(),
            noter={},
            metadata=metadata,
        )


def fetch_and_parse_arsredovisning(org_nummer: str, index: int = 0) -> Tuple[Arsredovisning, bytes, bytes]:
    """Hämta och parsa årsredovisning.
    
    Returnerar:
        Tuple med (Arsredovisning, xhtml_bytes, zip_bytes)
    """
    clean_nr = clean_org_nummer(org_nummer)
    
    dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
    dokument = dok_data.get("dokument", [])
    
    if not dokument:
        raise Exception("Inga årsredovisningar hittades")
    
    if index >= len(dokument):
        raise Exception(f"Index {index} finns inte. Det finns {len(dokument)} årsredovisningar.")
    
    dok = dokument[index]
    dok_id = dok.get("dokumentId")
    
    logger.info(f"Hämtar årsredovisning {index+1}/{len(dokument)} för {format_org_nummer(clean_nr)}")
    
    zip_bytes = download_document_bytes(dok_id)
    
    xhtml_content = None
    xhtml_filename = None
    with zipfile.ZipFile(BytesIO(zip_bytes)) as zf:
        for name in zf.namelist():
            if name.lower().endswith(('.xhtml', '.html', '.xml')):
                xhtml_content = zf.read(name).decode('utf-8')
                xhtml_filename = name
                break
    
    if not xhtml_content:
        raise Exception("Ingen XHTML-fil hittades i ZIP-arkivet")
    
    parser = IXBRLParser(xhtml_content)
    return parser.parse_full(), xhtml_content.encode('utf-8'), zip_bytes


# =============================================================================
# Export-funktioner
# =============================================================================

def export_to_json(data: Any) -> str:
    if hasattr(data, '__dataclass_fields__'):
        data = asdict(data)
    return json.dumps(data, indent=2, ensure_ascii=False, default=str)


def export_to_csv(data: Dict[str, Any], filename: str = None) -> str:
    output = StringIO()
    writer = csv.writer(output, delimiter=';')
    
    writer.writerow(['Nyckeltal', 'Värde', 'Enhet'])
    
    labels = {
        'nettoomsattning': ('Nettoomsättning', 'SEK'),
        'resultat_efter_finansiella': ('Resultat efter finansiella poster', 'SEK'),
        'arets_resultat': ('Årets resultat', 'SEK'),
        'eget_kapital': ('Eget kapital', 'SEK'),
        'balansomslutning': ('Balansomslutning', 'SEK'),
        'soliditet': ('Soliditet', '%'),
        'vinstmarginal': ('Vinstmarginal', '%'),
        'roe': ('Avkastning på eget kapital (ROE)', '%'),
        'antal_anstallda': ('Antal anställda', 'st'),
    }
    
    for key, (label, unit) in labels.items():
        value = data.get(key)
        if value is not None:
            writer.writerow([label, value, unit])
    
    csv_content = output.getvalue()
    
    if filename:
        filepath = ensure_output_dir() / filename
        with open(filepath, 'w', encoding='utf-8-sig') as f:
            f.write(csv_content)
        return str(filepath)
    
    return csv_content


def export_to_excel(arsredovisning: Arsredovisning, filename: str = None) -> str:
    if not EXCEL_AVAILABLE:
        return handle_error(ErrorCode.EXPORT_ERROR, "Excel-export ej tillgänglig", reason="openpyxl saknas")
    
    wb = openpyxl.Workbook()
    
    header_font = Font(bold=True, size=12)
    title_font = Font(bold=True, size=14)
    money_format = '#,##0'
    
    ws = wb.active
    ws.title = "Översikt"
    
    ws['A1'] = arsredovisning.foretag_namn
    ws['A1'].font = title_font
    ws['A2'] = f"Org.nr: {format_org_nummer(arsredovisning.org_nummer)}"
    ws['A3'] = f"Räkenskapsår: {arsredovisning.rakenskapsar_start} - {arsredovisning.rakenskapsar_slut}"
    
    ws['A5'] = 'Nyckeltal'
    ws['A5'].font = header_font
    
    row = 6
    nyckeltal = asdict(arsredovisning.nyckeltal)
    labels = {
        'nettoomsattning': 'Nettoomsättning',
        'resultat_efter_finansiella': 'Resultat efter finansiella poster',
        'arets_resultat': 'Årets resultat',
        'eget_kapital': 'Eget kapital',
        'balansomslutning': 'Balansomslutning',
        'soliditet': 'Soliditet (%)',
        'vinstmarginal': 'Vinstmarginal (%)',
        'roe': 'ROE (%)',
        'antal_anstallda': 'Antal anställda',
    }
    
    for key, label in labels.items():
        value = nyckeltal.get(key)
        if value is not None:
            ws[f'A{row}'] = label
            ws[f'B{row}'] = value
            if key not in ('soliditet', 'vinstmarginal', 'roe', 'antal_anstallda'):
                ws[f'B{row}'].number_format = money_format
            row += 1
    
    ws2 = wb.create_sheet("Personer")
    ws2['A1'] = 'Förnamn'
    ws2['B1'] = 'Efternamn'
    ws2['C1'] = 'Roll'
    for cell in ws2[1]:
        cell.font = header_font
    
    for i, person in enumerate(arsredovisning.personer, 2):
        ws2[f'A{i}'] = person.fornamn
        ws2[f'B{i}'] = person.efternamn
        ws2[f'C{i}'] = person.roll
    
    if not filename:
        clean_name = re.sub(r'[^\w\s-]', '', arsredovisning.foretag_namn)
        filename = f"{clean_name}_{arsredovisning.rakenskapsar_slut[:4]}.xlsx"
    
    filepath = ensure_output_dir() / filename
    wb.save(filepath)
    logger.info(f"Excel exporterad till: {filepath}")
    return str(filepath)


def export_to_pdf(arsredovisning: Arsredovisning, filename: str = None) -> str:
    if not PDF_AVAILABLE:
        return handle_error(ErrorCode.EXPORT_ERROR, "PDF-export ej tillgänglig", reason="weasyprint saknas")
    
    nyckeltal = arsredovisning.nyckeltal
    
    personer_html = ""
    for p in arsredovisning.personer:
        personer_html += f"<tr><td>{p.fornamn}</td><td>{p.efternamn}</td><td>{p.roll}</td></tr>"
    
    def fmt(val):
        return f"{val:,}" if val else "-"
    
    html_content = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="utf-8">
        <style>
            body {{ font-family: 'Helvetica Neue', Arial, sans-serif; margin: 40px; color: #333; }}
            h1 {{ color: #1a365d; border-bottom: 2px solid #1a365d; padding-bottom: 10px; }}
            h2 {{ color: #2c5282; margin-top: 30px; }}
            .info {{ background: #f7fafc; padding: 15px; border-radius: 5px; margin: 20px 0; }}
            table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
            th, td {{ padding: 10px; text-align: left; border-bottom: 1px solid #e2e8f0; }}
            th {{ background: #edf2f7; font-weight: bold; }}
            .number {{ text-align: right; font-family: monospace; }}
            .highlight {{ background: #ebf8ff; }}
        </style>
    </head>
    <body>
        <h1>{arsredovisning.foretag_namn}</h1>
        <div class="info">
            <strong>Organisationsnummer:</strong> {format_org_nummer(arsredovisning.org_nummer)}<br>
            <strong>Räkenskapsår:</strong> {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}
        </div>
        
        <h2>Nyckeltal</h2>
        <table>
            <tr><th>Nyckeltal</th><th class="number">Belopp (SEK)</th></tr>
            <tr><td>Nettoomsättning</td><td class="number">{fmt(nyckeltal.nettoomsattning)}</td></tr>
            <tr><td>Resultat efter finansiella poster</td><td class="number">{fmt(nyckeltal.resultat_efter_finansiella)}</td></tr>
            <tr class="highlight"><td><strong>Årets resultat</strong></td><td class="number"><strong>{fmt(nyckeltal.arets_resultat)}</strong></td></tr>
            <tr><td>Eget kapital</td><td class="number">{fmt(nyckeltal.eget_kapital)}</td></tr>
            <tr><td>Soliditet</td><td class="number">{nyckeltal.soliditet or '-'} %</td></tr>
            <tr><td>Vinstmarginal</td><td class="number">{nyckeltal.vinstmarginal or '-'} %</td></tr>
            <tr><td>ROE</td><td class="number">{nyckeltal.roe or '-'} %</td></tr>
        </table>
        
        <h2>Personer</h2>
        <table>
            <tr><th>Förnamn</th><th>Efternamn</th><th>Roll</th></tr>
            {personer_html}
        </table>
        
        <div style="margin-top: 40px; font-size: 11px; color: #718096;">
            Genererad: {datetime.now().strftime('%Y-%m-%d %H:%M')} | Källa: Bolagsverket
        </div>
    </body>
    </html>
    """
    
    if not filename:
        clean_name = re.sub(r'[^\w\s-]', '', arsredovisning.foretag_namn)
        filename = f"{clean_name}_{arsredovisning.rakenskapsar_slut[:4]}.pdf"
    
    filepath = ensure_output_dir() / filename
    HTML(string=html_content).write_pdf(filepath)
    logger.info(f"PDF exporterad till: {filepath}")
    return str(filepath)


def export_to_markdown(arsredovisning: Arsredovisning) -> str:
    nyckeltal = arsredovisning.nyckeltal
    
    lines = [
        f"# {arsredovisning.foretag_namn}",
        f"",
        f"**Organisationsnummer:** {format_org_nummer(arsredovisning.org_nummer)}  ",
        f"**Räkenskapsår:** {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}",
        f"",
        f"## Nyckeltal",
        f"",
        f"| Nyckeltal | Belopp |",
        f"|-----------|--------|",
    ]
    
    if nyckeltal.nettoomsattning:
        lines.append(f"| Nettoomsättning | {nyckeltal.nettoomsattning:,} SEK |")
    if nyckeltal.resultat_efter_finansiella:
        lines.append(f"| Resultat efter finansiella poster | {nyckeltal.resultat_efter_finansiella:,} SEK |")
    if nyckeltal.arets_resultat:
        lines.append(f"| **Årets resultat** | **{nyckeltal.arets_resultat:,} SEK** |")
    if nyckeltal.eget_kapital:
        lines.append(f"| Eget kapital | {nyckeltal.eget_kapital:,} SEK |")
    if nyckeltal.soliditet:
        lines.append(f"| Soliditet | {nyckeltal.soliditet} % |")
    if nyckeltal.vinstmarginal:
        lines.append(f"| Vinstmarginal | {nyckeltal.vinstmarginal} % |")
    if nyckeltal.roe:
        lines.append(f"| ROE | {nyckeltal.roe} % |")
    
    lines.extend([
        f"",
        f"## Personer",
        f"",
        f"| Namn | Roll |",
        f"|------|------|",
    ])
    
    for p in arsredovisning.personer:
        lines.append(f"| {p.fullnamn} | {p.roll} |")

    return "\n".join(lines)


def export_to_docx(arsredovisning: Arsredovisning, filename: str = None) -> str:
    """Exportera till Word-dokument (.docx)."""
    if not DOCX_AVAILABLE:
        return handle_error(ErrorCode.EXPORT_ERROR, "Word-export ej tillgänglig", reason="python-docx saknas")

    doc = Document()

    # Titel
    title = doc.add_heading(arsredovisning.foretag_namn, 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Företagsinfo
    doc.add_paragraph()
    info_para = doc.add_paragraph()
    info_para.add_run("Organisationsnummer: ").bold = True
    info_para.add_run(format_org_nummer(arsredovisning.org_nummer))
    info_para.add_run("\n")
    info_para.add_run("Räkenskapsår: ").bold = True
    info_para.add_run(f"{arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}")

    # Nyckeltal-sektion
    doc.add_heading("Nyckeltal", level=1)

    nyckeltal = arsredovisning.nyckeltal
    table = doc.add_table(rows=1, cols=2)
    table.style = 'Table Grid'

    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Nyckeltal'
    hdr_cells[1].text = 'Värde'
    for cell in hdr_cells:
        cell.paragraphs[0].runs[0].bold = True

    def fmt_money(val):
        return f"{val:,} SEK" if val else "-"

    def fmt_pct(val):
        return f"{val} %" if val is not None else "-"

    data = [
        ("Nettoomsättning", fmt_money(nyckeltal.nettoomsattning)),
        ("Resultat efter finansiella poster", fmt_money(nyckeltal.resultat_efter_finansiella)),
        ("Årets resultat", fmt_money(nyckeltal.arets_resultat)),
        ("Eget kapital", fmt_money(nyckeltal.eget_kapital)),
        ("Balansomslutning", fmt_money(nyckeltal.balansomslutning)),
        ("Soliditet", fmt_pct(nyckeltal.soliditet)),
        ("Vinstmarginal", fmt_pct(nyckeltal.vinstmarginal)),
        ("ROE", fmt_pct(nyckeltal.roe)),
        ("Antal anställda", str(nyckeltal.antal_anstallda) if nyckeltal.antal_anstallda else "-"),
    ]

    for label, value in data:
        row = table.add_row().cells
        row[0].text = label
        row[1].text = value

    # Personer-sektion
    doc.add_heading("Styrelse och ledning", level=1)

    if arsredovisning.personer:
        person_table = doc.add_table(rows=1, cols=2)
        person_table.style = 'Table Grid'

        hdr = person_table.rows[0].cells
        hdr[0].text = 'Namn'
        hdr[1].text = 'Roll'
        for cell in hdr:
            cell.paragraphs[0].runs[0].bold = True

        for person in arsredovisning.personer:
            row = person_table.add_row().cells
            row[0].text = person.fullnamn
            row[1].text = person.roll
    else:
        doc.add_paragraph("Inga personer registrerade i årsredovisningen.")

    # Footer
    doc.add_paragraph()
    footer = doc.add_paragraph()
    footer.add_run(f"Genererad: {datetime.now().strftime('%Y-%m-%d %H:%M')} | ").italic = True
    footer.add_run("Källa: Bolagsverket").italic = True

    # Spara
    if not filename:
        clean_name = re.sub(r'[^\w\s-]', '', arsredovisning.foretag_namn).strip().replace(' ', '_')
        year = arsredovisning.rakenskapsar_slut[:4] if arsredovisning.rakenskapsar_slut else "unknown"
        filename = f"{clean_name}_{year}_rapport.docx"

    filepath = ensure_output_dir() / filename
    doc.save(filepath)
    logger.info(f"Word-dokument exporterat till: {filepath}")
    return str(filepath)


def export_to_pptx(arsredovisning: Arsredovisning, filename: str = None) -> str:
    """Exportera till PowerPoint-presentation (.pptx)."""
    if not PPTX_AVAILABLE:
        return handle_error(ErrorCode.EXPORT_ERROR, "PowerPoint-export ej tillgänglig", reason="python-pptx saknas")

    prs = Presentation()
    prs.slide_width = PptxInches(13.333)  # 16:9 format
    prs.slide_height = PptxInches(7.5)

    nyckeltal = arsredovisning.nyckeltal

    # Slide 1: Titel
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Titel-box
    title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(2.5), PptxInches(12.333), PptxInches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = arsredovisning.foretag_namn
    p.font.size = PptxPt(44)
    p.font.bold = True
    p.font.color.rgb = PptxRGBColor(26, 54, 93)  # Mörkblå
    p.alignment = PP_ALIGN.CENTER

    # Underrubrik
    subtitle_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(4), PptxInches(12.333), PptxInches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Årsredovisning {arsredovisning.rakenskapsar_slut[:4]}"
    p.font.size = PptxPt(24)
    p.font.color.rgb = PptxRGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Org.nr
    org_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(5), PptxInches(12.333), PptxInches(0.5))
    tf = org_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Org.nr: {format_org_nummer(arsredovisning.org_nummer)}"
    p.font.size = PptxPt(18)
    p.font.color.rgb = PptxRGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Nyckeltal
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Rubrik
    header_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12.333), PptxInches(0.8))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Finansiella nyckeltal"
    p.font.size = PptxPt(32)
    p.font.bold = True
    p.font.color.rgb = PptxRGBColor(26, 54, 93)

    def fmt_money(val):
        if val is None:
            return "-"
        if abs(val) >= 1_000_000:
            return f"{val/1_000_000:.1f} MSEK"
        elif abs(val) >= 1_000:
            return f"{val/1_000:.0f} TSEK"
        return f"{val:,} SEK"

    # Nyckeltal i rutnät
    metrics = [
        ("Omsättning", fmt_money(nyckeltal.nettoomsattning)),
        ("Årets resultat", fmt_money(nyckeltal.arets_resultat)),
        ("Eget kapital", fmt_money(nyckeltal.eget_kapital)),
        ("Soliditet", f"{nyckeltal.soliditet}%" if nyckeltal.soliditet else "-"),
        ("Vinstmarginal", f"{nyckeltal.vinstmarginal}%" if nyckeltal.vinstmarginal else "-"),
        ("Anställda", str(nyckeltal.antal_anstallda) if nyckeltal.antal_anstallda else "-"),
    ]

    for i, (label, value) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = PptxInches(0.8 + col * 4.2)
        y = PptxInches(1.5 + row * 2.5)

        # Värde (stort)
        value_box = slide.shapes.add_textbox(x, y, PptxInches(3.8), PptxInches(1))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = PptxPt(36)
        p.font.bold = True
        p.font.color.rgb = PptxRGBColor(44, 82, 130) if not value.startswith("-") else PptxRGBColor(200, 50, 50)

        # Label (litet)
        label_box = slide.shapes.add_textbox(x, PptxInches(y.inches + 0.9), PptxInches(3.8), PptxInches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = PptxPt(16)
        p.font.color.rgb = PptxRGBColor(100, 100, 100)

    # Slide 3: Styrelse
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12.333), PptxInches(0.8))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Styrelse och ledning"
    p.font.size = PptxPt(32)
    p.font.bold = True
    p.font.color.rgb = PptxRGBColor(26, 54, 93)

    if arsredovisning.personer:
        y_pos = 1.3
        for person in arsredovisning.personer[:10]:  # Max 10 personer
            person_box = slide.shapes.add_textbox(PptxInches(1), PptxInches(y_pos), PptxInches(11), PptxInches(0.5))
            tf = person_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = person.fullnamn
            run.font.size = PptxPt(20)
            run.font.bold = True
            run = p.add_run()
            run.text = f"  –  {person.roll}"
            run.font.size = PptxPt(18)
            run.font.color.rgb = PptxRGBColor(100, 100, 100)
            y_pos += 0.55
    else:
        no_person_box = slide.shapes.add_textbox(PptxInches(1), PptxInches(2), PptxInches(11), PptxInches(1))
        tf = no_person_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Inga personer registrerade"
        p.font.size = PptxPt(18)
        p.font.italic = True

    # Slide 4: Källa
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    source_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(3), PptxInches(12.333), PptxInches(1.5))
    tf = source_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Källa: Bolagsverket"
    p.font.size = PptxPt(24)
    p.font.color.rgb = PptxRGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = f"Genererad: {datetime.now().strftime('%Y-%m-%d')}"
    p.font.size = PptxPt(16)
    p.font.color.rgb = PptxRGBColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER

    # Spara
    if not filename:
        clean_name = re.sub(r'[^\w\s-]', '', arsredovisning.foretag_namn).strip().replace(' ', '_')
        year = arsredovisning.rakenskapsar_slut[:4] if arsredovisning.rakenskapsar_slut else "unknown"
        filename = f"{clean_name}_{year}_presentation.pptx"

    filepath = ensure_output_dir() / filename
    prs.save(filepath)
    logger.info(f"PowerPoint exporterad till: {filepath}")
    return str(filepath)


# =============================================================================
# FÖRBÄTTRING #6: Pydantic-modeller med förbättrade inputSchema
# =============================================================================

class ResponseFormat(str, Enum):
    MARKDOWN = "markdown"
    JSON = "json"
    CSV = "csv"
    EXCEL = "excel"
    PDF = "pdf"
    DOCX = "docx"    # Word-dokument
    PPTX = "pptx"    # PowerPoint-presentation
    XHTML = "xhtml"  # Original iXBRL-fil
    ZIP = "zip"      # Original ZIP från Bolagsverket


class OrgNummerInput(BaseModel):
    """Input för organisationsnummer."""
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(
        min_length=10,
        max_length=13,
        description="10-siffrigt organisationsnummer (t.ex. 5567671267 eller 556767-1267)",
        json_schema_extra={"examples": ["5567671267", "556767-1267"]}
    )


class CompanyInfoInput(BaseModel):
    """Input för företagsinformation."""
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(
        min_length=10,
        description="10-siffrigt organisationsnummer"
    )
    response_format: ResponseFormat = Field(
        default=ResponseFormat.MARKDOWN,
        description="Svarsformat: markdown, json, csv, excel eller pdf"
    )


class FinansiellDataInput(BaseModel):
    """Input för finansiell data."""
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(min_length=10, description="Organisationsnummer")
    index: int = Field(
        default=0,
        ge=0,
        le=10,
        description="Vilken årsredovisning (0=senaste, 1=näst senaste)"
    )
    response_format: ResponseFormat = Field(default=ResponseFormat.MARKDOWN)


class BatchInput(BaseModel):
    """Input för batch-sökning."""
    model_config = ConfigDict(extra="forbid")
    org_nummer_lista: List[str] = Field(
        min_length=1,
        max_length=20,
        description="Lista med organisationsnummer att söka (max 20)"
    )


class ExportInput(BaseModel):
    """Input för export."""
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(min_length=10)
    index: int = Field(default=0, ge=0)
    format: ResponseFormat = Field(
        default=ResponseFormat.PDF,
        description="Exportformat: pdf, excel, csv, json eller markdown"
    )
    filename: Optional[str] = Field(
        default=None,
        description="Valfritt filnamn (genereras automatiskt om ej angivet)"
    )


class ArendenInput(BaseModel):
    """Input för ärenden/cases-sökning."""
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(
        min_length=10,
        description="10-siffrigt organisationsnummer"
    )
    from_datum: Optional[str] = Field(
        default=None,
        description="Startdatum för sökning (YYYY-MM-DD). Default: 1 år bakåt"
    )
    to_datum: Optional[str] = Field(
        default=None,
        description="Slutdatum för sökning (YYYY-MM-DD). Default: idag"
    )


# =============================================================================
# FÖRBÄTTRING #2 & #3: Resources (passiv data via URI-schema)
# =============================================================================

@mcp.resource("bolagsverket://company/{org_nummer}")
def resource_company(org_nummer: str) -> str:
    """
    Företagsinformation som resurs.
    
    URI: bolagsverket://company/{org_nummer}
    Exempel: bolagsverket://company/5567671267
    """
    try:
        logger.info(f"Resource request: company/{org_nummer}")
        valid, clean_nr = validate_org_nummer(org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr, org_nummer=org_nummer)
        
        info = fetch_company_info(clean_nr)
        return export_to_json(info)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


@mcp.resource("bolagsverket://financials/{org_nummer}")
def resource_financials(org_nummer: str) -> str:
    """
    Finansiell data (nyckeltal) som resurs.
    
    URI: bolagsverket://financials/{org_nummer}
    """
    try:
        logger.info(f"Resource request: financials/{org_nummer}")
        arsred, _, _ = fetch_and_parse_arsredovisning(org_nummer, 0)
        return export_to_json(arsred.nyckeltal)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


@mcp.resource("bolagsverket://people/{org_nummer}")
def resource_people(org_nummer: str) -> str:
    """
    Personer kopplade till företaget som resurs.
    
    URI: bolagsverket://people/{org_nummer}
    """
    try:
        logger.info(f"Resource request: people/{org_nummer}")
        arsred, _, _ = fetch_and_parse_arsredovisning(org_nummer, 0)
        return export_to_json([asdict(p) for p in arsred.personer])
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


@mcp.resource("bolagsverket://annual-reports/{org_nummer}")
def resource_annual_reports_list(org_nummer: str) -> str:
    """
    Lista tillgängliga årsredovisningar som resurs.
    
    URI: bolagsverket://annual-reports/{org_nummer}
    """
    try:
        logger.info(f"Resource request: annual-reports/{org_nummer}")
        clean_nr = clean_org_nummer(org_nummer)
        dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
        dokument = dok_data.get("dokument", [])
        
        result = []
        for i, dok in enumerate(dokument):
            result.append({
                "index": i,
                "dokument_id": dok.get("dokumentId"),
                "period_fran": dok.get("rakenskapsperiod", {}).get("fran"),
                "period_till": dok.get("rakenskapsperiod", {}).get("till"),
                "inlamningsdatum": dok.get("inlamningsdatum"),
            })
        
        return export_to_json({"org_nummer": format_org_nummer(clean_nr), "arsredovisningar": result})
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=org_nummer)


# =============================================================================
# FÖRBÄTTRING #4: Prompts (fördefinierade arbetsflöden)
# =============================================================================

@mcp.prompt("due-diligence")
def prompt_due_diligence(org_nummer: str) -> str:
    """
    Due diligence-analys av ett företag.
    
    Utför en komplett genomlysning med:
    - Grundläggande företagsinformation
    - Finansiell analys med nyckeltal
    - Identifiering av nyckelpersoner
    - Riskbedömning
    """
    return f"""Utför en due diligence-analys av företaget med organisationsnummer {org_nummer}.

## Steg 1: Grundläggande information
Använd verktyget `bolagsverket_get_basic_info` för att hämta:
- Företagsnamn och organisationsform
- Registreringsdatum och status
- Säte och adress
- Verksamhetsbeskrivning
- SNI-koder (bransch)

## Steg 2: Finansiell analys
Använd verktyget `bolagsverket_get_nyckeltal` för att analysera:
- Omsättning och resultat
- Soliditet och likviditet
- Vinstmarginal och ROE
- Jämför med föregående år om möjligt

## Steg 3: Nyckelpersoner
Använd verktyget `bolagsverket_get_styrelse` för att identifiera:
- VD och företrädare
- Styrelseledamöter
- Revisorer

## Steg 4: Trendanalys
Använd verktyget `bolagsverket_get_trends` för att:
- Analysera omsättningstillväxt
- Identifiera resultatutveckling
- Bedöma finansiell stabilitet över tid

## Steg 5: Sammanfattning
Ge en övergripande bedömning med:
- Styrkor
- Svagheter/risker
- Rekommendation
"""


@mcp.prompt("compare-companies")
def prompt_compare_companies(org_nummer_1: str, org_nummer_2: str) -> str:
    """
    Jämför två företag.
    
    Skapar en side-by-side jämförelse av nyckeltal och egenskaper.
    """
    return f"""Jämför följande två företag:

**Företag 1:** {org_nummer_1}
**Företag 2:** {org_nummer_2}

## Steg 1: Hämta grundinformation
Använd `bolagsverket_get_basic_info` för båda företagen.

## Steg 2: Hämta nyckeltal
Använd `bolagsverket_get_nyckeltal` för båda företagen.

## Steg 3: Skapa jämförelsetabell
Presentera en tabell med följande jämförelsepunkter:
- Omsättning
- Årets resultat
- Eget kapital
- Soliditet
- Vinstmarginal
- Antal anställda
- Organisationsform
- Bransch (SNI)

## Steg 4: Analys
- Vilket företag är större?
- Vilket har bättre lönsamhet?
- Vilket har starkare finansiell ställning?
- Eventuella risker eller fördelar
"""


@mcp.prompt("person-network")
def prompt_person_network(org_nummer: str) -> str:
    """
    Analysera personkopplingar för ett företag.
    
    Identifierar nyckelpersoner och deras roller.
    """
    return f"""Analysera nyckelpersoner kopplade till företaget {org_nummer}.

## Steg 1: Identifiera personer
Använd `bolagsverket_get_styrelse` för att hämta:
- VD
- Styrelseledamöter
- Revisorer

## Steg 2: Analysera roller
För varje person, notera:
- Fullständigt namn
- Roll i företaget
- Om de har flera roller

## Steg 3: Sammanfattning
Presentera en översikt av företagets ledning och governance-struktur.
"""


@mcp.prompt("export-report")
def prompt_export_report(org_nummer: str, format: str = "pdf") -> str:
    """
    Exportera en komplett företagsrapport.
    """
    return f"""Skapa och exportera en komplett rapport för företaget {org_nummer}.

## Steg 1: Samla data
Hämta all relevant information med:
- `bolagsverket_get_basic_info`
- `bolagsverket_get_nyckeltal`
- `bolagsverket_get_styrelse`

## Steg 2: Exportera
Använd `bolagsverket_export` med format="{format}" för att skapa en nedladdningsbar fil.

## Steg 3: Presentera
Ge användaren länk till den exporterade filen.
"""


# =============================================================================
# FÖRBÄTTRING #9: Granulära verktyg (Single Responsibility)
# =============================================================================

@mcp.tool()
def bolagsverket_check_status() -> str:
    """
    Kontrollera om Bolagsverkets API är tillgängligt.
    
    Returnerar API:ets status.
    """
    try:
        logger.info("Kontrollerar API-status...")
        token = token_manager.get_token()
        with httpx.Client(timeout=30.0) as client:
            response = client.get(f"{BASE_URL}/isalive", headers={"Authorization": f"Bearer {token}"})
        if response.status_code == 200:
            return "✅ Bolagsverkets API är tillgängligt!"
        return f"⚠️ API svarade med status {response.status_code}"
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_get_basic_info(params: OrgNummerInput) -> str:
    """
    Hämta grundläggande företagsinformation.
    
    Returnerar:
    - Företagsnamn
    - Organisationsform (AB, HB, etc.)
    - Juridisk form
    - Registreringsdatum
    - Status (aktiv/avregistrerad)
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result, org_nummer=params.org_nummer)
        
        info = fetch_company_info(result)
        
        lines = [
            f"# {info.namn}",
            f"",
            f"**Organisationsnummer:** {info.org_nummer}",
            f"**Organisationsform:** {info.organisationsform}",
        ]
        
        if info.juridisk_form:
            lines.append(f"**Juridisk form:** {info.juridisk_form}")
        
        lines.append(f"**Registreringsdatum:** {info.registreringsdatum}")
        lines.append(f"**Status:** {info.status}")
        
        if info.avregistreringsdatum:
            lines.append(f"**Avregistreringsdatum:** {info.avregistreringsdatum[:10]}")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e), org_nummer=params.org_nummer)


@mcp.tool()
def bolagsverket_get_address(params: OrgNummerInput) -> str:
    """
    Hämta företagets postadress.
    
    Returnerar:
    - Utdelningsadress
    - Postnummer
    - Postort
    - Säte (län)
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        info = fetch_company_info(result)
        
        lines = [
            f"# Adress: {info.namn}",
            f"",
        ]
        
        if info.adress.get('utdelningsadress'):
            lines.append(info.adress['utdelningsadress'])
        
        lines.append(f"{info.adress.get('postnummer', '')} {info.adress.get('postort', '')}")
        
        if info.sate:
            lines.append(f"**Säte:** {info.sate}")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_get_verksamhet(params: OrgNummerInput) -> str:
    """
    Hämta företagets verksamhetsbeskrivning och branschkoder.
    
    Returnerar:
    - Verksamhetsbeskrivning
    - SNI-koder med klartext
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        info = fetch_company_info(result)
        
        lines = [
            f"# Verksamhet: {info.namn}",
            f"",
        ]
        
        if info.verksamhet:
            lines.append("## Beskrivning")
            lines.append(info.verksamhet.strip())
            lines.append("")
        
        if info.sni_koder:
            lines.append("## SNI-koder (bransch)")
            for sni in info.sni_koder:
                lines.append(f"- **{sni['kod']}**: {sni['klartext']}")
        else:
            lines.append("*Inga SNI-koder registrerade*")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_get_nyckeltal(params: FinansiellDataInput) -> str:
    """
    Hämta finansiella nyckeltal från årsredovisning.
    
    Returnerar:
    - Nettoomsättning
    - Resultat efter finansiella poster
    - Årets resultat
    - Eget kapital
    - Soliditet
    - Vinstmarginal (beräknad)
    - ROE (beräknad)
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, _, _ = fetch_and_parse_arsredovisning(result, params.index)
        
        if params.response_format == ResponseFormat.JSON:
            return export_to_json(arsredovisning.nyckeltal)
        
        nyckeltal = arsredovisning.nyckeltal
        lines = [
            f"# Nyckeltal: {arsredovisning.foretag_namn}",
            f"**Räkenskapsår:** {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}",
            "",
            "| Nyckeltal | Värde |",
            "|-----------|------:|"
        ]
        
        data = [
            ("Nettoomsättning", nyckeltal.nettoomsattning, "SEK"),
            ("Resultat efter fin. poster", nyckeltal.resultat_efter_finansiella, "SEK"),
            ("Årets resultat", nyckeltal.arets_resultat, "SEK"),
            ("Eget kapital", nyckeltal.eget_kapital, "SEK"),
            ("Balansomslutning", nyckeltal.balansomslutning, "SEK"),
            ("Soliditet", nyckeltal.soliditet, "%"),
            ("Vinstmarginal", nyckeltal.vinstmarginal, "%"),
            ("ROE", nyckeltal.roe, "%"),
            ("Antal anställda", nyckeltal.antal_anstallda, "st"),
        ]
        
        for label, value, unit in data:
            if value is not None:
                if isinstance(value, int) and unit == "SEK":
                    lines.append(f"| {label} | {value:,} {unit} |")
                else:
                    lines.append(f"| {label} | {value} {unit} |")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.ANNUAL_REPORT_NOT_FOUND, str(e), org_nummer=params.org_nummer)


@mcp.tool()
def bolagsverket_get_styrelse(params: OrgNummerInput) -> str:
    """
    Hämta styrelse, VD och revisorer från årsredovisning.
    
    Returnerar:
    - VD (Verkställande direktör)
    - Styrelseledamöter
    - Revisorer
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, _, _ = fetch_and_parse_arsredovisning(result, 0)
        
        lines = [
            f"# Personer: {arsredovisning.foretag_namn}",
            f"**Org.nr:** {format_org_nummer(arsredovisning.org_nummer)}",
            f"**Källa:** Årsredovisning {arsredovisning.rakenskapsar_slut[:4]}",
            "",
            "| Namn | Roll |",
            "|------|------|",
        ]
        
        for person in arsredovisning.personer:
            lines.append(f"| {person.fullnamn} | {person.roll} |")
        
        if not arsredovisning.personer:
            lines.append("| *Inga personer hittades* | - |")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.ANNUAL_REPORT_NOT_FOUND, str(e), org_nummer=params.org_nummer)


@mcp.tool()
def bolagsverket_get_trends(params: OrgNummerInput) -> str:
    """
    Flerårsöversikt med trendanalys (upp till 4 år).
    
    Returnerar:
    - Nyckeltal per år
    - Omsättningstillväxt (%)
    - Resultatutveckling
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        clean_nr = clean_org_nummer(result)
        
        dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
        dokument = dok_data.get("dokument", [])
        
        if not dokument:
            return handle_error(ErrorCode.ANNUAL_REPORT_NOT_FOUND, "Inga årsredovisningar hittades")
        
        dok_id = dokument[0].get("dokumentId")
        zip_bytes = download_document_bytes(dok_id)
        
        xhtml_content = None
        with zipfile.ZipFile(BytesIO(zip_bytes)) as zf:
            for name in zf.namelist():
                if name.lower().endswith(('.xhtml', '.html')):
                    xhtml_content = zf.read(name).decode('utf-8')
                    break
        
        if not xhtml_content:
            return handle_error(ErrorCode.PARSE_ERROR, "Ingen XHTML hittades")
        
        parser = IXBRLParser(xhtml_content)
        metadata = parser.get_metadata()
        oversikt = parser.get_flerarsoversikt()
        
        lines = [
            f"# Flerårsöversikt: {metadata['foretag_namn']}",
            f"**Org.nr:** {format_org_nummer(metadata['org_nummer'])}",
            "",
        ]
        
        periods = sorted(oversikt.keys())
        if not periods:
            return handle_error(ErrorCode.PARSE_ERROR, "Kunde inte extrahera flerårsdata")
        
        header = "| Nyckeltal |"
        separator = "|-----------|"
        for p in periods:
            year = f"År {p[-1]}"
            header += f" {year} |"
            separator += "------:|"
        lines.append(header)
        lines.append(separator)
        
        metrics = [
            ('nettoomsattning', 'Nettoomsättning', True),
            ('resultat_efter_finansiella', 'Resultat', True),
            ('arets_resultat', 'Årets resultat', True),
            ('eget_kapital', 'Eget kapital', True),
            ('soliditet', 'Soliditet (%)', False),
        ]
        
        for key, label, is_money in metrics:
            row = f"| {label} |"
            for p in periods:
                val = getattr(oversikt[p], key, None)
                if val is not None:
                    if is_money:
                        row += f" {val:,} |"
                    else:
                        row += f" {val} |"
                else:
                    row += " - |"
            lines.append(row)
        
        if len(periods) >= 2:
            lines.append("")
            lines.append("## Tillväxt")
            
            first = oversikt[periods[-1]]
            last = oversikt[periods[0]]
            
            if first.nettoomsattning and last.nettoomsattning and first.nettoomsattning > 0:
                tillvaxt = ((last.nettoomsattning - first.nettoomsattning) / first.nettoomsattning) * 100
                lines.append(f"- **Omsättningstillväxt:** {tillvaxt:.1f}%")
            
            if first.arets_resultat and last.arets_resultat:
                diff = last.arets_resultat - first.arets_resultat
                lines.append(f"- **Resultatförändring:** {diff:+,} SEK")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_batch_lookup(params: BatchInput) -> str:
    """
    Sök information om flera företag samtidigt (max 20).
    
    Returnerar sammanställd tabell med:
    - Företagsnamn
    - Status
    - Omsättning
    - Resultat
    """
    try:
        results = []
        errors = []
        
        for org_nr in params.org_nummer_lista:
            try:
                clean_nr = clean_org_nummer(org_nr)
                
                org_data = make_api_request("POST", "/organisationer", {"identitetsbeteckning": clean_nr})
                orgs = org_data.get("organisationer", [])
                
                if not orgs:
                    errors.append(f"{org_nr}: Hittades ej")
                    continue
                
                org = orgs[0]
                namn = org.get("organisationsnamn", {}).get("organisationsnamnLista", [{}])[0].get("namn", "Okänt")
                form = org.get("organisationsform", {}).get("klartext", "-")
                
                avreg = org.get("avregistreradOrganisation", {})
                status = "Avregistrerad" if avreg and avreg.get("avregistreringsdatum") else "Aktiv"
                
                nyckeltal = None
                try:
                    arsred, _, _ = fetch_and_parse_arsredovisning(org_nr, 0)
                    nyckeltal = arsred.nyckeltal
                except:
                    pass
                
                results.append({
                    'org_nr': format_org_nummer(clean_nr),
                    'namn': namn,
                    'form': form,
                    'status': status,
                    'nyckeltal': nyckeltal
                })
                
            except Exception as e:
                errors.append(f"{org_nr}: {str(e)}")
        
        lines = [
            f"# Batch-sökning ({len(results)} företag)",
            "",
            "| Org.nr | Företag | Form | Status | Omsättning | Resultat |",
            "|--------|---------|------|--------|------------|----------|",
        ]
        
        for r in results:
            oms = f"{r['nyckeltal'].nettoomsattning:,}" if r['nyckeltal'] and r['nyckeltal'].nettoomsattning else "-"
            res = f"{r['nyckeltal'].arets_resultat:,}" if r['nyckeltal'] and r['nyckeltal'].arets_resultat else "-"
            lines.append(f"| {r['org_nr']} | {r['namn'][:25]} | {r['form'][:10]} | {r['status']} | {oms} | {res} |")
        
        if errors:
            lines.append("")
            lines.append("## Fel")
            for e in errors:
                lines.append(f"- {e}")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_export(params: ExportInput) -> str:
    """
    Exportera årsredovisningsdata till fil.

    Tillgängliga format:
    - pdf: Snygg PDF-rapport
    - excel: Excel med formatering
    - docx: Professionellt Word-dokument
    - pptx: PowerPoint-presentation (4 slides)
    - csv: CSV för import
    - json: Strukturerad JSON
    - markdown: Markdown-text
    - xhtml: Original iXBRL-fil från Bolagsverket
    - zip: Original ZIP-arkiv från Bolagsverket

    Filer sparas i ~/Downloads/bolagsverket/
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, xhtml_bytes, zip_bytes = fetch_and_parse_arsredovisning(result, params.index)
        
        # Generera filnamn
        clean_name = re.sub(r'[^\w\s-]', '', arsredovisning.foretag_namn).strip()
        year = arsredovisning.rakenskapsar_slut[:4] if arsredovisning.rakenskapsar_slut else "unknown"
        
        if params.format == ResponseFormat.ZIP:
            filename = params.filename or f"{clean_name}_{year}_arsredovisning.zip"
            filepath = ensure_output_dir() / filename
            with open(filepath, 'wb') as f:
                f.write(zip_bytes)
            logger.info(f"ZIP exporterad till: {filepath}")
            return f"✅ Original ZIP exporterad till: {filepath}"
        
        elif params.format == ResponseFormat.XHTML:
            filename = params.filename or f"{clean_name}_{year}_arsredovisning.xhtml"
            filepath = ensure_output_dir() / filename
            with open(filepath, 'wb') as f:
                f.write(xhtml_bytes)
            logger.info(f"XHTML exporterad till: {filepath}")
            return f"✅ Original XHTML (iXBRL) exporterad till: {filepath}"
        
        elif params.format == ResponseFormat.PDF:
            filepath = export_to_pdf(arsredovisning, params.filename)
            return f"✅ PDF exporterad till: {filepath}"
        
        elif params.format == ResponseFormat.EXCEL:
            filepath = export_to_excel(arsredovisning, params.filename)
            return f"✅ Excel exporterad till: {filepath}"
        
        elif params.format == ResponseFormat.CSV:
            filename = params.filename or f"{arsredovisning.foretag_namn}.csv"
            filepath = export_to_csv(asdict(arsredovisning.nyckeltal), filename)
            return f"✅ CSV exporterad till: {filepath}"

        elif params.format == ResponseFormat.DOCX:
            filepath = export_to_docx(arsredovisning, params.filename)
            return f"✅ Word-dokument exporterat till: {filepath}"

        elif params.format == ResponseFormat.PPTX:
            filepath = export_to_pptx(arsredovisning, params.filename)
            return f"✅ PowerPoint-presentation exporterad till: {filepath}"

        elif params.format == ResponseFormat.JSON:
            return export_to_json(arsredovisning)

        else:
            return export_to_markdown(arsredovisning)
        
    except Exception as e:
        return handle_error(ErrorCode.EXPORT_ERROR, str(e), format=params.format.value)


@mcp.tool()
def bolagsverket_list_arsredovisningar(params: OrgNummerInput) -> str:
    """
    Lista tillgängliga årsredovisningar för ett företag.
    
    Returnerar:
    - Räkenskapsperiod
    - Inlämningsdatum
    - Dokument-ID
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        clean_nr = clean_org_nummer(result)
        dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
        dokument = dok_data.get("dokument", [])
        
        if not dokument:
            return f"Inga årsredovisningar hittades för {format_org_nummer(clean_nr)}"
        
        lines = [
            f"# Årsredovisningar: {format_org_nummer(clean_nr)}",
            "",
            "| # | Period | Inlämnad |",
            "|---|--------|----------|",
        ]
        
        for i, dok in enumerate(dokument):
            period = dok.get("rakenskapsperiod", {})
            fran = period.get("fran", "-")
            till = period.get("till", "-")
            inlamnad = dok.get("inlamningsdatum", "-")
            lines.append(f"| {i} | {fran} – {till} | {inlamnad} |")
        
        lines.append("")
        lines.append(f"*Använd index 0-{len(dokument)-1} för att hämta specifik årsredovisning*")
        
        return "\n".join(lines)
    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


class DownloadInput(BaseModel):
    """Input för nedladdning av original-årsredovisning."""
    model_config = ConfigDict(extra="forbid")
    org_nummer: str = Field(min_length=10, description="Organisationsnummer")
    index: int = Field(default=0, ge=0, description="Vilken årsredovisning (0=senaste)")
    format: str = Field(
        default="zip",
        description="Format: 'zip' för original ZIP-arkiv, 'xhtml' för iXBRL-fil"
    )
    destination: Optional[str] = Field(
        default=None,
        description="Valfri destination (t.ex. ~/Desktop). Standard: ~/Downloads/bolagsverket/"
    )


@mcp.tool()
def bolagsverket_download_original(params: DownloadInput) -> str:
    """
    Ladda ner original-årsredovisning från Bolagsverket.
    
    Detta verktyg sparar den RIKTIGA årsredovisningen som Bolagsverket tillhandahåller,
    inte en genererad sammanfattning.
    
    Format:
    - zip: Original ZIP-arkiv som innehåller iXBRL-filen
    - xhtml: Extraherad iXBRL/XHTML-fil (kan öppnas i webbläsare)
    
    Returnerar sökväg till den sparade filen.
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)
        
        arsredovisning, xhtml_bytes, zip_bytes = fetch_and_parse_arsredovisning(result, params.index)
        
        # Generera filnamn
        clean_name = re.sub(r'[^\w\s-]', '', arsredovisning.foretag_namn).strip().replace(' ', '_')
        year = arsredovisning.rakenskapsar_slut[:4] if arsredovisning.rakenskapsar_slut else "unknown"
        
        # Bestäm destination
        if params.destination:
            dest_path = Path(params.destination).expanduser()
            dest_path.mkdir(parents=True, exist_ok=True)
        else:
            dest_path = ensure_output_dir()
        
        if params.format.lower() == "zip":
            filename = f"{clean_name}_{year}_arsredovisning.zip"
            filepath = dest_path / filename
            with open(filepath, 'wb') as f:
                f.write(zip_bytes)
            size_kb = len(zip_bytes) / 1024
            logger.info(f"Original ZIP sparad: {filepath} ({size_kb:.1f} KB)")
            return f"✅ Original årsredovisning (ZIP) sparad:\n📁 {filepath}\n📊 Storlek: {size_kb:.1f} KB\n📅 Räkenskapsår: {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}"
        
        elif params.format.lower() == "xhtml":
            filename = f"{clean_name}_{year}_arsredovisning.xhtml"
            filepath = dest_path / filename
            with open(filepath, 'wb') as f:
                f.write(xhtml_bytes)
            size_kb = len(xhtml_bytes) / 1024
            logger.info(f"Original XHTML sparad: {filepath} ({size_kb:.1f} KB)")
            return f"✅ Original årsredovisning (iXBRL/XHTML) sparad:\n📁 {filepath}\n📊 Storlek: {size_kb:.1f} KB\n📅 Räkenskapsår: {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}\n💡 Tips: Öppna filen i en webbläsare för att se årsredovisningen"
        
        else:
            return handle_error(ErrorCode.INVALID_INPUT, f"Okänt format: {params.format}. Använd 'zip' eller 'xhtml'.")
        
    except Exception as e:
        return handle_error(ErrorCode.EXPORT_ERROR, str(e))


# =============================================================================
# NYA VERKTYG: Djupanalys och jämförelse
# =============================================================================

class RiskLevel(str, Enum):
    CRITICAL = "KRITISK"
    HIGH = "HÖG"
    MEDIUM = "MEDIUM"
    LOW = "LÅG"
    INFO = "INFO"


@dataclass
class RiskFlag:
    """En riskindikator med allvarlighetsgrad."""
    level: RiskLevel
    category: str
    description: str
    value: Optional[str] = None
    recommendation: Optional[str] = None


def analyze_risks(arsredovisning: Arsredovisning, trends: Optional[Dict] = None) -> List[RiskFlag]:
    """Analysera risker baserat på nyckeltal och trender."""
    flags = []
    nyckeltal = arsredovisning.nyckeltal
    balans = arsredovisning.balansrakning

    # === KRITISKA FLAGGOR ===

    # Negativt eget kapital
    if nyckeltal.eget_kapital and nyckeltal.eget_kapital < 0:
        flags.append(RiskFlag(
            level=RiskLevel.CRITICAL,
            category="Kapitalstruktur",
            description="Negativt eget kapital",
            value=f"{nyckeltal.eget_kapital:,} SEK",
            recommendation="Bolaget kan behöva kapitaltillskott eller kontrollbalansräkning"
        ))

    # Kontrollbalansräkning (EK < 50% av aktiekapital)
    aktiekapital = balans.get('eget_kapital_skulder', {}).get('aktiekapital')
    if aktiekapital and nyckeltal.eget_kapital:
        if nyckeltal.eget_kapital < aktiekapital * 0.5:
            flags.append(RiskFlag(
                level=RiskLevel.CRITICAL,
                category="Kapitalstruktur",
                description="Eget kapital understiger 50% av aktiekapitalet",
                value=f"EK: {nyckeltal.eget_kapital:,} vs AK: {aktiekapital:,}",
                recommendation="Kontrollbalansräkning kan vara aktuell enligt ABL 25 kap"
            ))

    # Negativ soliditet
    if nyckeltal.soliditet and nyckeltal.soliditet < 0:
        flags.append(RiskFlag(
            level=RiskLevel.CRITICAL,
            category="Finansiell styrka",
            description="Negativ soliditet",
            value=f"{nyckeltal.soliditet}%",
            recommendation="Allvarlig finansiell obalans - överväg rekonstruktion"
        ))

    # === HÖGA FLAGGOR ===

    # Låg soliditet
    if nyckeltal.soliditet and 0 < nyckeltal.soliditet < 20:
        flags.append(RiskFlag(
            level=RiskLevel.HIGH,
            category="Finansiell styrka",
            description="Låg soliditet (under 20%)",
            value=f"{nyckeltal.soliditet}%",
            recommendation="Stärk eget kapital genom nyemission eller balanserade vinster"
        ))

    # Förlust
    if nyckeltal.arets_resultat and nyckeltal.arets_resultat < 0:
        flags.append(RiskFlag(
            level=RiskLevel.HIGH,
            category="Lönsamhet",
            description="Förlust under räkenskapsåret",
            value=f"{nyckeltal.arets_resultat:,} SEK",
            recommendation="Analysera kostnadsstruktur och intäktsmöjligheter"
        ))

    # Negativ vinstmarginal
    if nyckeltal.vinstmarginal and nyckeltal.vinstmarginal < -10:
        flags.append(RiskFlag(
            level=RiskLevel.HIGH,
            category="Lönsamhet",
            description="Kraftigt negativ vinstmarginal",
            value=f"{nyckeltal.vinstmarginal}%",
            recommendation="Omgående kostnadsöversyn behövs"
        ))

    # Hög skuldsättning
    if nyckeltal.eget_kapital and nyckeltal.balansomslutning:
        skulder = nyckeltal.balansomslutning - nyckeltal.eget_kapital
        if nyckeltal.eget_kapital > 0:
            skuldsattningsgrad = skulder / nyckeltal.eget_kapital
            if skuldsattningsgrad > 3:
                flags.append(RiskFlag(
                    level=RiskLevel.HIGH,
                    category="Skuldsättning",
                    description="Hög skuldsättningsgrad (över 3x)",
                    value=f"{skuldsattningsgrad:.1f}x",
                    recommendation="Överväg amortering eller kapitaltillskott"
                ))

    # === MEDIUM FLAGGOR ===

    # Låg vinstmarginal
    if nyckeltal.vinstmarginal and 0 < nyckeltal.vinstmarginal < 3:
        flags.append(RiskFlag(
            level=RiskLevel.MEDIUM,
            category="Lönsamhet",
            description="Låg vinstmarginal (under 3%)",
            value=f"{nyckeltal.vinstmarginal}%",
            recommendation="Överväg prissättning och kostnadseffektivisering"
        ))

    # Soliditet 20-30% (acceptabel men bör förbättras)
    if nyckeltal.soliditet and 20 <= nyckeltal.soliditet < 30:
        flags.append(RiskFlag(
            level=RiskLevel.MEDIUM,
            category="Finansiell styrka",
            description="Måttlig soliditet (20-30%)",
            value=f"{nyckeltal.soliditet}%",
            recommendation="Fortsätt stärka eget kapital"
        ))

    # === LÅG/INFO FLAGGOR ===

    # Inga anställda rapporterade
    if not nyckeltal.antal_anstallda:
        flags.append(RiskFlag(
            level=RiskLevel.INFO,
            category="Personal",
            description="Inga anställda rapporterade",
            recommendation="Kan vara korrekt för holdingbolag eller enmansföretag"
        ))

    # Trendbaserade flaggor (om trend-data finns)
    if trends:
        years = sorted(trends.keys())
        if len(years) >= 2:
            first_year = trends[years[-1]]
            last_year = trends[years[0]]

            # Fallande omsättning
            if first_year.nettoomsattning and last_year.nettoomsattning:
                if first_year.nettoomsattning > 0:
                    change = ((last_year.nettoomsattning - first_year.nettoomsattning) / first_year.nettoomsattning) * 100
                    if change < -20:
                        flags.append(RiskFlag(
                            level=RiskLevel.HIGH,
                            category="Tillväxt",
                            description="Fallande omsättning (över 20%)",
                            value=f"{change:.1f}%",
                            recommendation="Analysera marknadsposition och konkurrenskraft"
                        ))
                    elif change < -10:
                        flags.append(RiskFlag(
                            level=RiskLevel.MEDIUM,
                            category="Tillväxt",
                            description="Sjunkande omsättning (10-20%)",
                            value=f"{change:.1f}%",
                            recommendation="Bevaka marknadsutvecklingen"
                        ))

    return flags


@mcp.tool()
def bolagsverket_risk_analysis(params: OrgNummerInput) -> str:
    """
    Djupgående riskanalys av ett företag.

    Analyserar:
    - Kapitalstruktur (eget kapital, kontrollbalansräkning)
    - Finansiell styrka (soliditet, skuldsättning)
    - Lönsamhet (vinstmarginal, ROE)
    - Trender (om flerårsdata finns)

    Returnerar röda flaggor klassificerade efter allvarlighetsgrad:
    - KRITISK: Kräver omedelbar åtgärd
    - HÖG: Allvarlig risk
    - MEDIUM: Bör åtgärdas
    - LÅG: Bevaka
    - INFO: Information
    """
    try:
        valid, result = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, result)

        # Hämta årsredovisning
        arsredovisning, _, _ = fetch_and_parse_arsredovisning(result, 0)

        # Försök hämta flerårsdata för trendanalys
        trends = None
        try:
            clean_nr = clean_org_nummer(result)
            dok_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
            dokument = dok_data.get("dokument", [])

            if dokument:
                dok_id = dokument[0].get("dokumentId")
                zip_bytes = download_document_bytes(dok_id)

                with zipfile.ZipFile(BytesIO(zip_bytes)) as zf:
                    for name in zf.namelist():
                        if name.lower().endswith(('.xhtml', '.html')):
                            xhtml_content = zf.read(name).decode('utf-8')
                            parser = IXBRLParser(xhtml_content)
                            trends = parser.get_flerarsoversikt()
                            break
        except Exception as e:
            logger.warning(f"Kunde inte hämta trenddata: {e}")

        # Analysera risker
        flags = analyze_risks(arsredovisning, trends)
        nyckeltal = arsredovisning.nyckeltal

        # Bygg rapport
        lines = [
            f"# Riskanalys: {arsredovisning.foretag_namn}",
            f"**Org.nr:** {format_org_nummer(arsredovisning.org_nummer)}",
            f"**Räkenskapsår:** {arsredovisning.rakenskapsar_start} – {arsredovisning.rakenskapsar_slut}",
            "",
        ]

        # Sammanfattning
        critical_count = sum(1 for f in flags if f.level == RiskLevel.CRITICAL)
        high_count = sum(1 for f in flags if f.level == RiskLevel.HIGH)
        medium_count = sum(1 for f in flags if f.level == RiskLevel.MEDIUM)

        if critical_count > 0:
            lines.append(f"## ⛔ RISKNIVÅ: KRITISK")
            lines.append(f"*{critical_count} kritiska varningar kräver omedelbar uppmärksamhet*")
        elif high_count > 0:
            lines.append(f"## 🔴 RISKNIVÅ: HÖG")
            lines.append(f"*{high_count} höga varningar identifierade*")
        elif medium_count > 0:
            lines.append(f"## 🟡 RISKNIVÅ: MEDIUM")
            lines.append(f"*{medium_count} varningar bör åtgärdas*")
        else:
            lines.append(f"## 🟢 RISKNIVÅ: LÅG")
            lines.append(f"*Inga allvarliga risker identifierade*")

        lines.append("")

        # Nyckeltal
        lines.append("## Finansiell översikt")
        lines.append("")
        lines.append("| Nyckeltal | Värde |")
        lines.append("|-----------|------:|")

        if nyckeltal.nettoomsattning:
            lines.append(f"| Omsättning | {nyckeltal.nettoomsattning:,} SEK |")
        if nyckeltal.arets_resultat:
            lines.append(f"| Årets resultat | {nyckeltal.arets_resultat:,} SEK |")
        if nyckeltal.eget_kapital:
            lines.append(f"| Eget kapital | {nyckeltal.eget_kapital:,} SEK |")
        if nyckeltal.soliditet:
            lines.append(f"| Soliditet | {nyckeltal.soliditet}% |")
        if nyckeltal.vinstmarginal:
            lines.append(f"| Vinstmarginal | {nyckeltal.vinstmarginal}% |")

        lines.append("")

        # Flaggor grupperade efter nivå
        if flags:
            lines.append("## Identifierade risker")
            lines.append("")

            for level in [RiskLevel.CRITICAL, RiskLevel.HIGH, RiskLevel.MEDIUM, RiskLevel.LOW, RiskLevel.INFO]:
                level_flags = [f for f in flags if f.level == level]
                if level_flags:
                    level_emoji = {"KRITISK": "⛔", "HÖG": "🔴", "MEDIUM": "🟡", "LÅG": "🔵", "INFO": "ℹ️"}
                    lines.append(f"### {level_emoji.get(level.value, '')} {level.value}")
                    for flag in level_flags:
                        lines.append(f"- **{flag.category}:** {flag.description}")
                        if flag.value:
                            lines.append(f"  - *Värde:* {flag.value}")
                        if flag.recommendation:
                            lines.append(f"  - *Rekommendation:* {flag.recommendation}")
                    lines.append("")
        else:
            lines.append("## ✅ Inga risker identifierade")
            lines.append("Företaget visar inga uppenbara varningssignaler baserat på tillgänglig data.")

        return "\n".join(lines)

    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


class CompareInput(BaseModel):
    """Input för företagsjämförelse."""
    model_config = ConfigDict(extra="forbid")
    org_nummer_1: str = Field(min_length=10, description="Första företagets organisationsnummer")
    org_nummer_2: str = Field(min_length=10, description="Andra företagets organisationsnummer")


@mcp.tool()
def bolagsverket_compare(params: CompareInput) -> str:
    """
    Jämför två företag side-by-side.

    Jämför:
    - Storlek (omsättning, balansomslutning, anställda)
    - Lönsamhet (vinstmarginal, ROE)
    - Finansiell styrka (soliditet, skuldsättning)
    - Risknivå

    Returnerar en detaljerad jämförelsetabell.
    """
    try:
        # Validera båda org.nr
        valid1, clean_nr1 = validate_org_nummer(params.org_nummer_1)
        valid2, clean_nr2 = validate_org_nummer(params.org_nummer_2)

        if not valid1:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr1)
        if not valid2:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr2)

        # Hämta data för båda företagen
        try:
            arsred1, _, _ = fetch_and_parse_arsredovisning(clean_nr1, 0)
            info1 = fetch_company_info(clean_nr1)
        except Exception as e:
            return handle_error(ErrorCode.COMPANY_NOT_FOUND, f"Företag 1: {e}")

        try:
            arsred2, _, _ = fetch_and_parse_arsredovisning(clean_nr2, 0)
            info2 = fetch_company_info(clean_nr2)
        except Exception as e:
            return handle_error(ErrorCode.COMPANY_NOT_FOUND, f"Företag 2: {e}")

        n1 = arsred1.nyckeltal
        n2 = arsred2.nyckeltal

        # Beräkna skuldsättning
        def calc_skuldsattning(nyckeltal):
            if nyckeltal.eget_kapital and nyckeltal.balansomslutning and nyckeltal.eget_kapital > 0:
                skulder = nyckeltal.balansomslutning - nyckeltal.eget_kapital
                return round(skulder / nyckeltal.eget_kapital, 2)
            return None

        skuld1 = calc_skuldsattning(n1)
        skuld2 = calc_skuldsattning(n2)

        # Analysera risker
        flags1 = analyze_risks(arsred1)
        flags2 = analyze_risks(arsred2)

        critical1 = sum(1 for f in flags1 if f.level == RiskLevel.CRITICAL)
        high1 = sum(1 for f in flags1 if f.level == RiskLevel.HIGH)
        critical2 = sum(1 for f in flags2 if f.level == RiskLevel.CRITICAL)
        high2 = sum(1 for f in flags2 if f.level == RiskLevel.HIGH)

        risk_score1 = critical1 * 10 + high1 * 3
        risk_score2 = critical2 * 10 + high2 * 3

        def fmt(val, suffix=""):
            if val is None:
                return "-"
            if isinstance(val, int):
                return f"{val:,}{suffix}"
            return f"{val}{suffix}"

        def winner(val1, val2, higher_is_better=True):
            if val1 is None and val2 is None:
                return "", ""
            if val1 is None:
                return "", " ✓" if higher_is_better else ""
            if val2 is None:
                return " ✓" if higher_is_better else "", ""
            if val1 > val2:
                return " ✓" if higher_is_better else "", "" if higher_is_better else " ✓"
            elif val2 > val1:
                return "" if higher_is_better else " ✓", " ✓" if higher_is_better else ""
            return "", ""

        lines = [
            f"# Företagsjämförelse",
            "",
            f"| | **{arsred1.foretag_namn[:25]}** | **{arsred2.foretag_namn[:25]}** |",
            "|---|---:|---:|",
            f"| Org.nr | {format_org_nummer(arsred1.org_nummer)} | {format_org_nummer(arsred2.org_nummer)} |",
            f"| Bransch | {info1.sni_koder[0]['klartext'][:20] if info1.sni_koder else '-'} | {info2.sni_koder[0]['klartext'][:20] if info2.sni_koder else '-'} |",
            "",
            "## Storlek",
            "",
            f"| Nyckeltal | {arsred1.foretag_namn[:20]} | {arsred2.foretag_namn[:20]} |",
            "|---|---:|---:|",
        ]

        w1, w2 = winner(n1.nettoomsattning, n2.nettoomsattning)
        lines.append(f"| Omsättning | {fmt(n1.nettoomsattning)}{w1} | {fmt(n2.nettoomsattning)}{w2} |")

        w1, w2 = winner(n1.balansomslutning, n2.balansomslutning)
        lines.append(f"| Balansomslutning | {fmt(n1.balansomslutning)}{w1} | {fmt(n2.balansomslutning)}{w2} |")

        w1, w2 = winner(n1.antal_anstallda, n2.antal_anstallda)
        lines.append(f"| Anställda | {fmt(n1.antal_anstallda)}{w1} | {fmt(n2.antal_anstallda)}{w2} |")

        lines.extend([
            "",
            "## Lönsamhet",
            "",
            f"| Nyckeltal | {arsred1.foretag_namn[:20]} | {arsred2.foretag_namn[:20]} |",
            "|---|---:|---:|",
        ])

        w1, w2 = winner(n1.arets_resultat, n2.arets_resultat)
        lines.append(f"| Årets resultat | {fmt(n1.arets_resultat)}{w1} | {fmt(n2.arets_resultat)}{w2} |")

        w1, w2 = winner(n1.vinstmarginal, n2.vinstmarginal)
        lines.append(f"| Vinstmarginal | {fmt(n1.vinstmarginal, '%')}{w1} | {fmt(n2.vinstmarginal, '%')}{w2} |")

        w1, w2 = winner(n1.roe, n2.roe)
        lines.append(f"| ROE | {fmt(n1.roe, '%')}{w1} | {fmt(n2.roe, '%')}{w2} |")

        lines.extend([
            "",
            "## Finansiell styrka",
            "",
            f"| Nyckeltal | {arsred1.foretag_namn[:20]} | {arsred2.foretag_namn[:20]} |",
            "|---|---:|---:|",
        ])

        w1, w2 = winner(n1.eget_kapital, n2.eget_kapital)
        lines.append(f"| Eget kapital | {fmt(n1.eget_kapital)}{w1} | {fmt(n2.eget_kapital)}{w2} |")

        w1, w2 = winner(n1.soliditet, n2.soliditet)
        lines.append(f"| Soliditet | {fmt(n1.soliditet, '%')}{w1} | {fmt(n2.soliditet, '%')}{w2} |")

        w1, w2 = winner(skuld1, skuld2, higher_is_better=False)
        lines.append(f"| Skuldsättning | {fmt(skuld1, 'x')}{w1} | {fmt(skuld2, 'x')}{w2} |")

        lines.extend([
            "",
            "## Risknivå",
            "",
            f"| | {arsred1.foretag_namn[:20]} | {arsred2.foretag_namn[:20]} |",
            "|---|---:|---:|",
        ])

        w1, w2 = winner(risk_score1, risk_score2, higher_is_better=False)
        lines.append(f"| Kritiska flaggor | {critical1}{w1} | {critical2}{w2} |")
        lines.append(f"| Höga flaggor | {high1} | {high2} |")
        lines.append(f"| Risk-score | {risk_score1}{w1} | {risk_score2}{w2} |")

        # Sammanfattning
        lines.extend([
            "",
            "## Sammanfattning",
            "",
        ])

        # Storlek
        if n1.nettoomsattning and n2.nettoomsattning:
            if n1.nettoomsattning > n2.nettoomsattning:
                ratio = n1.nettoomsattning / n2.nettoomsattning
                lines.append(f"- **{arsred1.foretag_namn}** är {ratio:.1f}x större i omsättning")
            else:
                ratio = n2.nettoomsattning / n1.nettoomsattning
                lines.append(f"- **{arsred2.foretag_namn}** är {ratio:.1f}x större i omsättning")

        # Lönsamhet
        if n1.vinstmarginal and n2.vinstmarginal:
            if n1.vinstmarginal > n2.vinstmarginal:
                lines.append(f"- **{arsred1.foretag_namn}** har bättre vinstmarginal ({n1.vinstmarginal}% vs {n2.vinstmarginal}%)")
            else:
                lines.append(f"- **{arsred2.foretag_namn}** har bättre vinstmarginal ({n2.vinstmarginal}% vs {n1.vinstmarginal}%)")

        # Finansiell styrka
        if n1.soliditet and n2.soliditet:
            if n1.soliditet > n2.soliditet:
                lines.append(f"- **{arsred1.foretag_namn}** har starkare finansiell ställning (soliditet {n1.soliditet}% vs {n2.soliditet}%)")
            else:
                lines.append(f"- **{arsred2.foretag_namn}** har starkare finansiell ställning (soliditet {n2.soliditet}% vs {n1.soliditet}%)")

        # Risk
        if risk_score1 < risk_score2:
            lines.append(f"- **{arsred1.foretag_namn}** har lägre risknivå")
        elif risk_score2 < risk_score1:
            lines.append(f"- **{arsred2.foretag_namn}** har lägre risknivå")
        else:
            lines.append("- Båda företagen har liknande risknivå")

        return "\n".join(lines)

    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


# =============================================================================
# MVP FASE 1: Nya verktyg med hög journalistisk prioritet
# =============================================================================

@mcp.tool()
def bolagsverket_arenden(params: ArendenInput) -> str:
    """
    🔔 HÄNDELSEHISTORIK: Se alla registrerade händelser för ett bolag.

    Visar tidslinje med:
    - Alla inlämnade årsredovisningar
    - Registreringsdatum
    - Avregistrering (om tillämpligt)
    - Aktuell status

    JOURNALISTISK NYTTA: Snabb översikt över bolagets registreringshistorik.

    Args:
        org_nummer: Organisationsnummer
        from_datum: Startdatum (default: 5 år bakåt)
        to_datum: Slutdatum (default: idag)
    """
    try:
        valid, clean_nr = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr)

        # Hämta företagsinfo
        try:
            info = fetch_company_info(clean_nr)
        except Exception as e:
            return handle_error(ErrorCode.COMPANY_NOT_FOUND, str(e))

        # Hämta dokumentlista (årsredovisningar)
        try:
            doc_data = make_api_request("POST", "/dokumentlista", {"identitetsbeteckning": clean_nr})
            dokument = doc_data.get("dokument", [])
        except:
            dokument = []

        # Beräkna datumintervall för filtrering
        today = datetime.now()
        if params.to_datum:
            to_date = datetime.strptime(params.to_datum, "%Y-%m-%d")
        else:
            to_date = today

        if params.from_datum:
            from_date = datetime.strptime(params.from_datum, "%Y-%m-%d")
        else:
            from_date = today - timedelta(days=365 * 5)  # 5 år bakåt

        # Samla händelser
        events = []

        # Registreringsdatum
        if info.registreringsdatum:
            events.append({
                "datum": info.registreringsdatum,
                "typ": "🟢 REGISTRERING",
                "beskrivning": f"Bolaget registrerat som {info.organisationsform}"
            })

        # Avregistrering om sådan finns
        if info.avregistreringsdatum:
            events.append({
                "datum": info.avregistreringsdatum,
                "typ": "🔴 AVREGISTRERING",
                "beskrivning": "Bolaget avregistrerat"
            })

        # Årsredovisningar
        for doc in dokument:
            slut = doc.get("rapporteringsperiodTom", "")
            registrerad = doc.get("registreringstidpunkt", "")

            if slut:
                events.append({
                    "datum": registrerad or slut,
                    "typ": "📄 ÅRSREDOVISNING",
                    "beskrivning": f"Räkenskapsår t.o.m. {slut}"
                })

        # Sortera efter datum (nyast först)
        events.sort(key=lambda x: x.get("datum", ""), reverse=True)

        # Filtrera på datumintervall
        filtered_events = []
        for event in events:
            try:
                event_date = datetime.strptime(event["datum"][:10], "%Y-%m-%d")
                if from_date <= event_date <= to_date:
                    filtered_events.append(event)
            except:
                filtered_events.append(event)

        # Formatera output
        lines = [
            f"# 🔔 Händelser för {info.namn}",
            f"**Org.nr:** {format_org_nummer(clean_nr)}",
            f"**Status:** {'🔴 Avregistrerad' if info.avregistreringsdatum else '🟢 Aktiv'}",
            f"**Organisationsform:** {info.organisationsform}",
            f"**Registrerad:** {info.registreringsdatum}",
            ""
        ]

        if info.avregistreringsdatum:
            lines.append(f"⚠️ **VARNING:** Bolaget avregistrerades {info.avregistreringsdatum}")
            lines.append("")

        lines.append(f"## 📅 Tidslinje ({len(filtered_events)} händelser)")
        lines.append("")

        if not filtered_events:
            lines.append("*Inga händelser under vald period.*")
        else:
            for event in filtered_events[:20]:
                lines.append(f"- **{event['datum'][:10]}** {event['typ']}: {event['beskrivning']}")

            if len(filtered_events) > 20:
                lines.append(f"*...och {len(filtered_events) - 20} fler händelser*")

        # Sammanfattning
        num_arsred = sum(1 for e in events if "ÅRSREDOVISNING" in e["typ"])
        lines.append("")
        lines.append("---")
        lines.append(f"📊 **Sammanfattning:** {num_arsred} årsredovisningar registrerade.")

        if num_arsred == 0:
            lines.append("⚠️ **OBS:** Inga årsredovisningar hittades - nystartat eller undantaget bolag?")

        return "\n".join(lines)

    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_styrelse_ledning(params: OrgNummerInput) -> str:
    """
    👥 STYRELSE & LEDNING: Lista alla personer kopplade till ett bolag.

    Visar:
    - Styrelseledamöter
    - VD/Verkställande direktör
    - Revisorer
    - Suppleanter
    - Firmatecknare

    JOURNALISTISK NYTTA: Se vem som styr bolaget och kan citeras.

    Args:
        org_nummer: Organisationsnummer
    """
    try:
        valid, clean_nr = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr)

        # Hämta årsredovisning för att extrahera styrelse
        try:
            arsred, _, _ = fetch_and_parse_arsredovisning(clean_nr, 0)
            info = fetch_company_info(clean_nr)
        except Exception as e:
            return handle_error(ErrorCode.ANNUAL_REPORT_NOT_FOUND, str(e))

        personer = arsred.personer

        # Kategorisera personer efter roll
        kategorier = {
            "👔 VD / Verkställande direktör": [],
            "🏛️ Styrelseordförande": [],
            "📋 Styrelseledamöter": [],
            "📝 Suppleanter": [],
            "🔍 Revisorer": [],
            "✍️ Övriga": []
        }

        for person in personer:
            roll_lower = person.roll.lower() if person.roll else ""
            entry = f"**{person.fullnamn}** - {person.roll}"

            if "vd" in roll_lower or "verkställande" in roll_lower:
                kategorier["👔 VD / Verkställande direktör"].append(entry)
            elif "ordförande" in roll_lower:
                kategorier["🏛️ Styrelseordförande"].append(entry)
            elif "suppleant" in roll_lower:
                kategorier["📝 Suppleanter"].append(entry)
            elif "revisor" in roll_lower:
                kategorier["🔍 Revisorer"].append(entry)
            elif "ledamot" in roll_lower or "styrelse" in roll_lower:
                kategorier["📋 Styrelseledamöter"].append(entry)
            else:
                kategorier["✍️ Övriga"].append(entry)

        # Formatera output
        lines = [
            f"# 👥 Styrelse & Ledning",
            f"**{arsred.foretag_namn}** ({format_org_nummer(arsred.org_nummer)})",
            f"**Källa:** Årsredovisning {arsred.rakenskapsar_slut[:4]}",
            f"**Bransch:** {info.sni_koder[0]['klartext'] if info.sni_koder else 'Okänd'}",
            ""
        ]

        if not personer:
            lines.append("*Inga personer hittades i årsredovisningen.*")
            lines.append("")
            lines.append("💡 **Tips:** Prova `bolagsverket_get_styrelse` för funktionärer från bolagsregistret.")
        else:
            total = 0
            for kategori, personer_lista in kategorier.items():
                if personer_lista:
                    lines.append(f"## {kategori}")
                    for p in personer_lista:
                        lines.append(f"- {p}")
                    lines.append("")
                    total += len(personer_lista)

            lines.append("---")
            lines.append(f"📊 **Totalt:** {total} personer i styrelse och ledning.")

        return "\n".join(lines)

    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


@mcp.tool()
def bolagsverket_finansiell_snapshot(params: OrgNummerInput) -> str:
    """
    📊 FINANSIELL SAMMANFATTNING: Förstå företagets ekonomi på 30 sekunder.

    Visar kompakt finansiell bild:
    - Omsättning och resultat
    - Finansiell ställning (soliditet, skuldsättning)
    - Röda flaggor (varningar)
    - Jämförelse med föregående år

    JOURNALISTISK NYTTA: Snabbt förstå "hur går det för bolaget?"

    Args:
        org_nummer: Organisationsnummer
    """
    try:
        valid, clean_nr = validate_org_nummer(params.org_nummer)
        if not valid:
            return handle_error(ErrorCode.INVALID_INPUT, clean_nr)

        # Hämta senaste och föregående årsredovisning
        try:
            arsred, _, _ = fetch_and_parse_arsredovisning(clean_nr, 0)
            info = fetch_company_info(clean_nr)
        except Exception as e:
            return handle_error(ErrorCode.ANNUAL_REPORT_NOT_FOUND, str(e))

        n = arsred.nyckeltal

        # Försök hämta förra året för jämförelse
        prev_year = None
        try:
            prev_arsred, _, _ = fetch_and_parse_arsredovisning(clean_nr, 1)
            prev_year = prev_arsred.nyckeltal
        except:
            pass

        # Beräkna ytterligare nyckeltal
        skuldsattning = None
        if n.eget_kapital and n.balansomslutning and n.eget_kapital > 0:
            skulder = n.balansomslutning - n.eget_kapital
            skuldsattning = round(skulder / n.eget_kapital, 2)

        # Analysera risker
        flags = analyze_risks(arsred)
        critical = [f for f in flags if f.level == RiskLevel.CRITICAL]
        high = [f for f in flags if f.level == RiskLevel.HIGH]
        medium = [f for f in flags if f.level == RiskLevel.MEDIUM]

        # Bestäm övergripande status
        if critical:
            status_emoji = "🔴"
            status_text = "KRITISK RISK"
        elif high:
            status_emoji = "🟠"
            status_text = "FÖRHÖJD RISK"
        elif medium:
            status_emoji = "🟡"
            status_text = "NORMAL"
        else:
            status_emoji = "🟢"
            status_text = "STABIL"

        def fmt_sek(val, scale=1):
            if val is None:
                return "-"
            if abs(val) >= 1_000_000:
                return f"{val/1_000_000:.1f} MSEK"
            elif abs(val) >= 1_000:
                return f"{val/1_000:.0f} TSEK"
            return f"{val:,.0f} SEK"

        def trend_arrow(current, previous):
            if current is None or previous is None:
                return ""
            if current > previous * 1.1:
                return " ↗️"
            elif current < previous * 0.9:
                return " ↘️"
            return " →"

        # Formatera output
        lines = [
            f"# 📊 Finansiell Snapshot",
            f"**{arsred.foretag_namn}** ({format_org_nummer(arsred.org_nummer)})",
            f"**Räkenskapsår:** {arsred.rakenskapsar_start} – {arsred.rakenskapsar_slut}",
            f"**Bransch:** {info.sni_koder[0]['klartext'] if info.sni_koder else 'Okänd'}",
            "",
            f"## {status_emoji} Status: {status_text}",
            "",
            "## 💰 Resultat",
            "",
            "| Nyckeltal | Värde | Förändring |",
            "|---|---:|---:|",
        ]

        # Omsättning
        prev_oms = prev_year.nettoomsattning if prev_year else None
        trend = trend_arrow(n.nettoomsattning, prev_oms)
        change = ""
        if n.nettoomsattning and prev_oms and prev_oms > 0:
            pct = ((n.nettoomsattning - prev_oms) / prev_oms) * 100
            change = f"{pct:+.0f}%"
        lines.append(f"| Omsättning | {fmt_sek(n.nettoomsattning)}{trend} | {change} |")

        # Resultat
        prev_res = prev_year.arets_resultat if prev_year else None
        trend = trend_arrow(n.arets_resultat, prev_res)
        change = ""
        if n.arets_resultat and prev_res and prev_res != 0:
            if prev_res > 0:
                pct = ((n.arets_resultat - prev_res) / abs(prev_res)) * 100
                change = f"{pct:+.0f}%"
        lines.append(f"| Årets resultat | {fmt_sek(n.arets_resultat)}{trend} | {change} |")

        # Vinstmarginal
        lines.append(f"| Vinstmarginal | {n.vinstmarginal}% | |" if n.vinstmarginal else "| Vinstmarginal | - | |")

        lines.append("")
        lines.append("## 🏦 Finansiell Ställning")
        lines.append("")
        lines.append("| Nyckeltal | Värde | Bedömning |")
        lines.append("|---|---:|---|")

        # Eget kapital
        if n.eget_kapital:
            if n.eget_kapital < 0:
                bedomning = "⚠️ Negativt"
            else:
                bedomning = "✓"
            lines.append(f"| Eget kapital | {fmt_sek(n.eget_kapital)} | {bedomning} |")

        # Soliditet
        if n.soliditet:
            if n.soliditet < 20:
                bedomning = "⚠️ Svag (<20%)"
            elif n.soliditet < 30:
                bedomning = "⚡ Acceptabel"
            else:
                bedomning = "✓ God"
            lines.append(f"| Soliditet | {n.soliditet}% | {bedomning} |")

        # Skuldsättning
        if skuldsattning is not None:
            if skuldsattning > 3:
                bedomning = "⚠️ Hög"
            elif skuldsattning > 2:
                bedomning = "⚡ Förhöjd"
            else:
                bedomning = "✓ Normal"
            lines.append(f"| Skuldsättning | {skuldsattning}x | {bedomning} |")

        lines.append(f"| Anställda | {n.antal_anstallda or 0} st | |")

        # Varningar
        if critical or high:
            lines.append("")
            lines.append("## ⚠️ Varningar")
            for flag in critical:
                lines.append(f"- 🔴 **{flag.category}**: {flag.description}")
            for flag in high[:3]:
                lines.append(f"- 🟠 **{flag.category}**: {flag.description}")

        # Quick facts
        lines.append("")
        lines.append("---")
        lines.append("**Quick Facts:**")
        facts = []
        if n.nettoomsattning:
            if n.nettoomsattning > 100_000_000:
                facts.append("Storföretag (>100 MSEK)")
            elif n.nettoomsattning > 10_000_000:
                facts.append("Medelstort företag")
            else:
                facts.append("Småföretag")

        if n.soliditet and n.soliditet > 50:
            facts.append("Stark finansiell ställning")

        if n.arets_resultat and n.arets_resultat > 0:
            facts.append("Lönsamt")
        elif n.arets_resultat and n.arets_resultat < 0:
            facts.append("Förlust")

        lines.append(" • ".join(facts) if facts else "*Ingen sammanfattning tillgänglig*")

        return "\n".join(lines)

    except Exception as e:
        return handle_error(ErrorCode.API_ERROR, str(e))


# =============================================================================
# Huvudprogram
# =============================================================================

if __name__ == "__main__":
    logger.info("Startar Bolagsverket MCP Server v3.2...")
    mcp.run(transport="stdio")
